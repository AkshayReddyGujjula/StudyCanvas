"""
Converts .docx and .pptx files to high-quality, text-selectable PDFs.

Conversion priority chain (best quality first)
-----------------------------------------------
1. Microsoft Office COM automation (Windows only, requires Word / PowerPoint)
   → Produces output byte-for-byte identical to "File → Save As → PDF".
2. LibreOffice headless subprocess (any OS, requires LibreOffice installed)
   → Near-identical output; free and cross-platform.
3. ReportLab pure-Python fallback (works everywhere, no external tools needed)
   → Used on Vercel serverless where neither Office nor LibreOffice is available.
   → Supports: headings, bold/italic/underline, font sizes, colours, text
     alignment, bulleted/numbered lists, tables, images, and clickable hyperlinks.

The first converter that succeeds is used.  All failures are logged at DEBUG
level so the chain degrades gracefully without surfacing internal error details.
"""

import io
import logging
import os
import shutil
import subprocess
import sys
import tempfile
import uuid
from typing import Optional

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# Shared helpers
# ─────────────────────────────────────────────────────────────────────────────

def _count_pdf_pages(pdf_bytes: bytes) -> int:
    """Return the page count of an in-memory PDF (uses pypdf, always available)."""
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))
        return max(len(reader.pages), 1)
    except Exception:
        return 1


# ─────────────────────────────────────────────────────────────────────────────
# Converter 1 — Microsoft Office COM (Windows)
# ─────────────────────────────────────────────────────────────────────────────

def _try_word_com(file_path: str) -> Optional[bytes]:
    """
    Convert a .docx file using Word COM automation via docx2pdf.
    Returns PDF bytes on success, None if Word is unavailable or conversion fails.
    """
    if sys.platform != "win32":
        return None
    try:
        import docx2pdf
        with tempfile.TemporaryDirectory() as tmpdir:
            out_path = os.path.join(tmpdir, "output.pdf")
            docx2pdf.convert(os.path.abspath(file_path), out_path)
            if os.path.exists(out_path):
                with open(out_path, "rb") as f:
                    return f.read()
    except ImportError:
        logger.debug("docx2pdf not installed, skipping Word COM path")
    except Exception as exc:
        logger.debug("Word COM conversion failed: %s", exc)
    return None


def _try_powerpoint_com(file_path: str) -> Optional[bytes]:
    """
    Convert a .pptx file using PowerPoint COM automation (pywin32 / win32com).
    Returns PDF bytes on success, None if PowerPoint is unavailable.

    THREADING NOTE: comtypes.client.CreateObject silently fails when called from
    asyncio.to_thread() worker threads because comtypes initialises COM once on
    the main thread (pure-Python, no DllMain).  win32com.client is backed by the
    native pythoncom DLL which auto-initialises COM per-thread via DLL_THREAD_ATTACH,
    exactly as docx2pdf does for Word.  We also call pythoncom.CoInitialize()
    explicitly to guarantee STA mode (required by PowerPoint).
    """
    if sys.platform != "win32":
        return None
    try:
        import pythoncom          # part of pywin32 — installed as docx2pdf dependency
        import win32com.client as win32
    except ImportError:
        logger.warning("pywin32 not installed — skipping PowerPoint COM path")
        return None

    pythoncom.CoInitialize()      # initialise STA for this thread-pool thread
    try:
        abs_in = os.path.abspath(file_path)
        with tempfile.TemporaryDirectory() as tmpdir:
            abs_out = os.path.join(tmpdir, "output.pdf")
            ppt_app = win32.Dispatch("PowerPoint.Application")
            ppt_app.Visible = 1   # required — SaveAs silently fails when hidden
            try:
                presentation = ppt_app.Presentations.Open(
                    abs_in,
                    ReadOnly=-1,  # VARIANT_BOOL True = -1  (not Python True = 1)
                    Untitled=0,
                    WithWindow=0,
                )
                try:
                    # ppSaveAsPDF = 32  (Office 2010+)
                    presentation.SaveAs(os.path.abspath(abs_out), 32)
                finally:
                    presentation.Close()
            finally:
                ppt_app.Quit()
            if os.path.exists(abs_out):
                with open(abs_out, "rb") as f:
                    return f.read()
    except Exception as exc:
        logger.warning("PowerPoint COM conversion failed: %s", exc)
    finally:
        pythoncom.CoUninitialize()  # always release the STA apartment

    return None


# ─────────────────────────────────────────────────────────────────────────────
# Converter 2 — LibreOffice subprocess
# ─────────────────────────────────────────────────────────────────────────────

def _try_libreoffice(file_path: str) -> Optional[bytes]:
    """
    Convert any Office file to PDF using LibreOffice in headless mode.
    Returns PDF bytes on success, None if LibreOffice is not found or fails.

    A unique user-profile directory is created per conversion to allow
    concurrent invocations without lock-file conflicts.
    """
    # Locate the LibreOffice binary
    candidates: list[str] = ["soffice", "libreoffice"]
    if sys.platform == "win32":
        candidates += [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
    elif sys.platform == "darwin":
        candidates += ["/Applications/LibreOffice.app/Contents/MacOS/soffice"]

    lo_binary: Optional[str] = None
    for candidate in candidates:
        if shutil.which(candidate) or os.path.isfile(candidate):
            lo_binary = candidate
            break
    if not lo_binary:
        return None

    # Unique user profile so concurrent conversions don't collide
    profile_dir = os.path.join(tempfile.gettempdir(), f"lo_prof_{uuid.uuid4().hex}")
    os.makedirs(profile_dir, exist_ok=True)

    if sys.platform == "win32":
        profile_url = "file:///" + profile_dir.replace("\\", "/")
    else:
        profile_url = "file://" + profile_dir

    try:
        with tempfile.TemporaryDirectory() as out_dir:
            result = subprocess.run(
                [
                    lo_binary,
                    f"-env:UserInstallation={profile_url}",
                    "--headless",
                    "--norestore",
                    "--nofirststartwizard",
                    "--convert-to", "pdf",
                    "--outdir", out_dir,
                    file_path,
                ],
                timeout=120,
                capture_output=True,
                text=True,
            )

            if result.returncode != 0:
                logger.warning(
                    "LibreOffice exited with rc=%d: %s", result.returncode, result.stderr[:200]
                )
                return None

            # Locate the output PDF (LibreOffice preserves the stem name)
            stem = os.path.splitext(os.path.basename(file_path))[0]
            out_pdf = os.path.join(out_dir, stem + ".pdf")
            if not os.path.exists(out_pdf):
                pdfs = [f for f in os.listdir(out_dir) if f.endswith(".pdf")]
                if pdfs:
                    out_pdf = os.path.join(out_dir, pdfs[0])
                else:
                    logger.warning("LibreOffice produced no PDF output")
                    return None

            with open(out_pdf, "rb") as f:
                return f.read()

    except subprocess.TimeoutExpired:
        logger.warning("LibreOffice conversion timed out for %s", os.path.basename(file_path))
    except Exception as exc:
        logger.debug("LibreOffice conversion error: %s", exc)
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)

    return None


# ─────────────────────────────────────────────────────────────────────────────
# Converter 3 — ReportLab pure-Python fallback
# ─────────────────────────────────────────────────────────────────────────────
# Shared markup / colour helpers used by both DOCX and PPTX renderers.

def _esc(text: str) -> str:
    """Escape text for ReportLab XML markup."""
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _rgb_to_hex(rgb_color) -> Optional[str]:
    """
    Safely convert a python-pptx / python-docx RGBColor to a CSS hex string.

    RGBColor inherits from int and exposes .r / .g / .b — NOT index-based access.
    Using rgb[0] raises TypeError which was silently swallowed, dropping all colours.
    """
    if rgb_color is None:
        return None
    try:
        return f"#{rgb_color.r:02x}{rgb_color.g:02x}{rgb_color.b:02x}"
    except Exception:
        return None


# Alias map: pptx bg1/tx1/bg2/tx2 → real clrScheme key
_THEME_COLOR_ALIASES: dict[str, str] = {
    "bg1": "lt1", "tx1": "dk1", "bg2": "lt2", "tx2": "dk2",
}


def _build_theme_color_map(slide) -> dict[str, str]:
    """
    Return a dict mapping scheme color names → 6-hex strings (uppercase, no #),
    e.g. {'dk1': '000000', 'lt1': 'FFFFFF', 'accent1': '4472C4', ...}.

    Walks: slide → slide_layout → slide_master → theme part → <a:clrScheme>.
    Handles srgbClr, sysClr (uses lastClr), and scrgbClr (float percentages).
    Returns {} on any failure so callers degrade gracefully.
    """
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    try:
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        master = slide.slide_layout.slide_master
        theme_part = master.part.part_related_by(RT.THEME)
        clr_scheme = theme_part.element.find(f".//{{{A}}}clrScheme")
        if clr_scheme is None:
            return {}
        result: dict[str, str] = {}
        for elem in clr_scheme:
            tag = elem.tag.split("}")[-1]  # e.g. "dk1", "accent1"
            srgb = elem.find(f"{{{A}}}srgbClr")
            if srgb is not None:
                val = srgb.get("val", "")
                if val:
                    result[tag] = val.upper()
                continue
            sys_clr = elem.find(f"{{{A}}}sysClr")
            if sys_clr is not None:
                last = sys_clr.get("lastClr", "")
                if last:
                    result[tag] = last.upper()
                continue
            sc = elem.find(f"{{{A}}}scrgbClr")
            if sc is not None:
                try:
                    r = round(int(sc.get("r", "0")) * 255 / 100000)
                    g = round(int(sc.get("g", "0")) * 255 / 100000)
                    b = round(int(sc.get("b", "0")) * 255 / 100000)
                    result[tag] = f"{r:02X}{g:02X}{b:02X}"
                except Exception:
                    pass
        return result
    except Exception:
        return {}


def _apply_brightness(hex6: str, brightness: float) -> str:
    """
    Apply a pptx lumMod/lumOff brightness value (−1.0 … 1.0) to a 6-char hex color.
    Positive → tint (lighten); negative → shade (darken).
    """
    import colorsys
    try:
        r, g, b = int(hex6[0:2], 16) / 255, int(hex6[2:4], 16) / 255, int(hex6[4:6], 16) / 255
        h, lum, s = colorsys.rgb_to_hls(r, g, b)
        lum = lum + (1.0 - lum) * brightness if brightness > 0 else lum * (1.0 + brightness)
        lum = max(0.0, min(1.0, lum))
        nr, ng, nb = colorsys.hls_to_rgb(h, lum, s)
        return f"{round(nr * 255):02X}{round(ng * 255):02X}{round(nb * 255):02X}"
    except Exception:
        return hex6


def _resolve_color_hex(color_format, theme_map: dict[str, str]) -> Optional[str]:
    """
    Safely resolve a python-pptx / python-docx ColorFormat to '#rrggbb'.

    Handles:
    - MSO_COLOR_TYPE.RGB     → direct .rgb.r/.g/.b
    - MSO_COLOR_TYPE.SCHEME  → look up scheme key in theme_map + apply brightness
    - type is None           → not set at this level; return None (inherited)
    - anything else          → return None (rare: HSL, PRESET, SCRGB, SYSTEM)
    """
    try:
        from pptx.enum.dml import MSO_COLOR_TYPE
        t = color_format.type
    except Exception:
        return None

    if t is None:
        return None

    try:
        if t == MSO_COLOR_TYPE.RGB:
            rgb = color_format.rgb
            return f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"

        if t == MSO_COLOR_TYPE.SCHEME:
            xml_val = color_format.theme_color.xml_value          # e.g. "accent1"
            key = _THEME_COLOR_ALIASES.get(xml_val, xml_val)
            hex6 = theme_map.get(key)
            if hex6:
                try:
                    brightness = color_format.brightness          # float, 0 = no change
                    if brightness != 0.0:
                        hex6 = _apply_brightness(hex6, brightness)
                except Exception:
                    pass
                return "#" + hex6.lower()
    except Exception:
        pass

    return None


def _get_run_hyperlink(run, slide_part) -> Optional[str]:
    """
    Extract the external hyperlink URL from a python-pptx Run via its XML.

    python-pptx's run.hyperlink.address misses links stored on the rPr element
    (the most common case in modern PPTX files). This walks the XML directly:
      <a:rPr> → <a:hlinkClick r:id="rId5"/> → slide_part.rels[rId5].target_ref
    Returns the URL string or None.
    """
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    try:
        rPr = run._r.find(f"{{{A}}}rPr")
        if rPr is None:
            return None
        hlink = rPr.find(f"{{{A}}}hlinkClick")
        if hlink is None:
            return None
        r_id = hlink.get(f"{{{R}}}id")
        if not r_id:
            return None
        return slide_part.rels[r_id].target_ref
    except Exception:
        return None


# ── DOCX ReportLab renderer ──────────────────────────────────────────────────

def _docx_reportlab(file_path: str) -> tuple[bytes, int]:
    """
    Convert a .docx file to PDF using ReportLab Platypus.
    Preserves: headings, paragraph styles, bold / italic / underline / colour /
    font-size, text alignment, bulleted / numbered lists, tables, and hyperlinks.
    """
    from docx import Document
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak,
    )

    doc = Document(file_path)
    buffer = io.BytesIO()

    # Honour document page size / margins
    page_width, page_height = A4
    lm = rm = tm = bm = 72.0
    try:
        sec = doc.sections[0]
        if sec.page_width and sec.page_height:
            pw, ph = float(sec.page_width.pt), float(sec.page_height.pt)
            if pw > 0 and ph > 0:
                page_width, page_height = pw, ph
        if sec.left_margin:   lm = float(sec.left_margin.pt)
        if sec.right_margin:  rm = float(sec.right_margin.pt)
        if sec.top_margin:    tm = float(sec.top_margin.pt)
        if sec.bottom_margin: bm = float(sec.bottom_margin.pt)
    except Exception:
        pass

    pdf_doc = SimpleDocTemplate(
        buffer, pagesize=(page_width, page_height),
        leftMargin=lm, rightMargin=rm, topMargin=tm, bottomMargin=bm,
    )
    base = getSampleStyleSheet()

    W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    # ── alignment ──

    def _rl_align(align) -> int:
        try:
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            return {
                WD_ALIGN_PARAGRAPH.LEFT:    TA_LEFT,
                WD_ALIGN_PARAGRAPH.CENTER:  TA_CENTER,
                WD_ALIGN_PARAGRAPH.RIGHT:   TA_RIGHT,
                WD_ALIGN_PARAGRAPH.JUSTIFY: TA_JUSTIFY,
            }.get(align, TA_LEFT)
        except Exception:
            return TA_LEFT

    # ── single run → HTML ──

    def _run_elem_to_html(r_elem, para) -> str:
        """Convert one <w:r> XML element to ReportLab HTML markup."""
        from docx.text.run import Run as DocxRun
        try:
            run = DocxRun(r_elem, para)
        except Exception:
            return ""
        raw = run.text
        if not raw:
            return ""
        text = _esc(raw)

        open_tags: list[str] = []
        close_tags: list[str] = []

        if run.bold:
            open_tags.append("<b>");  close_tags.insert(0, "</b>")
        if run.italic:
            open_tags.append("<i>");  close_tags.insert(0, "</i>")
        if run.underline:
            open_tags.append("<u>");  close_tags.insert(0, "</u>")

        font_attrs = ""
        if run.font.size:
            try:
                font_attrs += f' size="{max(6, min(72, int(run.font.size.pt)))}"'
            except Exception:
                pass
        try:
            fc = run.font.color
            if fc and fc.type is not None:
                hex_c = _resolve_color_hex(fc, {})
                if hex_c:
                    font_attrs += f' color="{hex_c}"'
        except Exception:
            pass

        if font_attrs:
            open_tags.append(f"<font{font_attrs}>")
            close_tags.insert(0, "</font>")

        return "".join(open_tags) + text + "".join(close_tags)

    def _para_xml_to_html(para) -> str:
        """
        Walk paragraph XML to produce HTML, correctly handling:
        - normal <w:r> runs
        - <w:hyperlink> elements (extracts URL and wraps text in <link> tag)
        """
        parts: list[str] = []
        for elem in para._p:
            tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

            if tag == "hyperlink":
                # Resolve URL via document relationships
                r_id = elem.get(f"{{{R_NS}}}id")
                url = ""
                if r_id:
                    try:
                        url = para.part.rels[r_id].target_ref
                    except Exception:
                        pass

                # Collect HTML from inner runs
                inner_parts: list[str] = []
                for r_elem in elem.iter(f"{{{W_NS}}}r"):
                    inner_parts.append(_run_elem_to_html(r_elem, para))
                inner_html = "".join(inner_parts)

                if url and inner_html.strip():
                    safe_url = url.replace("&", "&amp;").replace('"', "&quot;")
                    parts.append(f'<link href="{safe_url}" color="blue">{inner_html}</link>')
                else:
                    parts.append(inner_html)

            elif tag == "r":
                parts.append(_run_elem_to_html(elem, para))

        return "".join(parts)

    def _para_style(para) -> ParagraphStyle:
        sname = (para.style.name if para.style else "") or "Normal"
        align = _rl_align(para.alignment)
        run_size: Optional[int] = None
        for r in para.runs:
            if r.font.size:
                try:
                    run_size = max(6, min(72, int(r.font.size.pt))); break
                except Exception:
                    pass
        if "Title" in sname:
            fs = run_size or 24
            return ParagraphStyle("_title", parent=base["Title"],
                                  fontSize=fs, leading=fs * 1.2, spaceAfter=8, alignment=align)
        for prefix, dfs, pname in [
            ("Heading 1", 18, "Heading1"), ("Heading 2", 15, "Heading2"),
            ("Heading 3", 13, "Heading3"), ("Heading 4", 12, "Heading4"),
            ("Heading 5", 11, "Heading4"), ("Heading 6", 10, "Heading4"),
        ]:
            if sname.startswith(prefix):
                fs = run_size or dfs
                return ParagraphStyle(prefix, parent=base[pname],
                                      fontSize=fs, leading=fs * 1.3, spaceAfter=6, alignment=align)
        fs = run_size or 11
        is_list = "List Bullet" in sname or "List Number" in sname
        return ParagraphStyle("_normal", parent=base["Normal"],
                              fontSize=fs, leading=fs * 1.3, spaceAfter=4,
                              leftIndent=18 if is_list else 0, alignment=align)

    def _para_to_flowables(para) -> list:
        html = _para_xml_to_html(para)
        sname = (para.style.name if para.style else "") or "Normal"
        if "List Bullet" in sname:
            html = "• " + html
        elif "List Number" in sname:
            html = "› " + html
        if not html.strip():
            return [Spacer(1, 4)]
        ps = _para_style(para)
        try:
            return [Paragraph(html, ps)]
        except Exception:
            return [Paragraph(_esc(para.text) or " ", ps)]

    def _table_to_flowables(table) -> list:
        data = []
        for row in table.rows:
            row_cells = []
            for cell in row.cells:
                cell_story = []
                for p in cell.paragraphs:
                    h = _para_xml_to_html(p)
                    if h.strip():
                        try:
                            cell_story.append(Paragraph(h, base["Normal"]))
                        except Exception:
                            cell_story.append(Paragraph(_esc(p.text) or " ", base["Normal"]))
                row_cells.append(cell_story if cell_story else [Paragraph(" ", base["Normal"])])
            if row_cells:
                data.append(row_cells)
        if not data:
            return []
        tbl = Table(data, hAlign="LEFT", repeatRows=1)
        tbl.setStyle(TableStyle([
            ("GRID",          (0, 0), (-1, -1), 0.5, colors.Color(0.7, 0.7, 0.7)),
            ("BACKGROUND",    (0, 0), (-1,  0), colors.Color(0.92, 0.92, 0.92)),
            ("VALIGN",        (0, 0), (-1, -1), "TOP"),
            ("TOPPADDING",    (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ("LEFTPADDING",   (0, 0), (-1, -1), 5),
            ("RIGHTPADDING",  (0, 0), (-1, -1), 5),
        ]))
        return [tbl, Spacer(1, 8)]

    story: list = []
    page_count_est = 1

    for block in doc.element.body:
        tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag

        if tag == "p":
            from docx.text.paragraph import Paragraph as DocxPara
            para = DocxPara(block, doc)
            for elem in block.iter():
                br_tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
                if br_tag == "br" and elem.get(f"{{{W_NS}}}type") == "page":
                    story.append(PageBreak())
                    page_count_est += 1
                    break
            story.extend(_para_to_flowables(para))

        elif tag == "tbl":
            from docx.table import Table as DocxTable
            story.extend(_table_to_flowables(DocxTable(block, doc)))

        elif tag == "sdt":
            from docx.text.paragraph import Paragraph as DocxPara
            for p_elem in block.iter(f"{{{W_NS}}}p"):
                story.extend(_para_to_flowables(DocxPara(p_elem, doc)))

    if not story:
        story.append(Paragraph("(empty document)", base["Normal"]))

    pdf_doc.build(story)
    return buffer.getvalue(), max(page_count_est, 1)


# ── PPTX ReportLab renderer ──────────────────────────────────────────────────

def _pptx_reportlab(file_path: str) -> tuple[bytes, int]:
    """
    Convert a .pptx file to PDF using ReportLab canvas (one page per slide).
    Text is selectable; layout mirrors original slide positions.
    Colour bug fixed: uses rgb.r / rgb.g / rgb.b (not rgb[0] / rgb[1] / rgb[2]).
    Hyperlinks: rendered as clickable blue <link> annotations.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from reportlab.pdfgen import canvas as rl_canvas
    from reportlab.platypus import Frame, Paragraph as RLPara, Spacer
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.utils import ImageReader
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT

    prs = Presentation(file_path)
    buffer = io.BytesIO()

    slide_w = float((prs.slide_width  or 9_144_000).pt)
    slide_h = float((prs.slide_height or 6_858_000).pt)

    c = rl_canvas.Canvas(buffer, pagesize=(slide_w, slide_h))

    def _rl_align(pp_align) -> int:
        return {PP_ALIGN.CENTER: TA_CENTER, PP_ALIGN.RIGHT: TA_RIGHT}.get(pp_align, TA_LEFT)

    def _draw_shape(shape, theme_map: dict, slide_part) -> None:
        try:
            x    = float(shape.left.pt   if shape.left   else 0)
            ytop = float(shape.top.pt    if shape.top    else 0)
            w    = float(shape.width.pt  if shape.width  else 0)
            h    = float(shape.height.pt if shape.height else 0)
        except Exception:
            return

        # ReportLab Y is from page bottom; PPTX Y is from slide top
        ybot = slide_h - ytop - h

        # ── Picture ──────────────────────────────────────────────────────────
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                c.drawImage(
                    ImageReader(io.BytesIO(shape.image.blob)),
                    x, ybot, width=w, height=h,
                    preserveAspectRatio=False, anchor="sw",
                )
            except Exception as exc:
                logger.debug("Skipping slide image: %s", exc)
            return

        # ── Group: recurse ───────────────────────────────────────────────────
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                for child in shape.shapes:
                    _draw_shape(child, theme_map, slide_part)
            except Exception:
                pass
            return

        if not shape.has_text_frame:
            return

        # ── Shape background fill (solid only) ───────────────────────────────
        try:
            fill = shape.fill
            if fill.type is not None:
                hex_c = _resolve_color_hex(fill.fore_color, theme_map)
                if hex_c:
                    c.saveState()
                    r, g, b = int(hex_c[1:3], 16), int(hex_c[3:5], 16), int(hex_c[5:7], 16)
                    c.setFillColorRGB(r / 255.0, g / 255.0, b / 255.0)
                    c.rect(x, ybot, w, h, fill=1, stroke=0)
                    c.restoreState()
        except Exception:
            pass

        # ── Text frame → ReportLab story inside a Frame ──────────────────────
        story: list = []
        for para in shape.text_frame.paragraphs:
            html_parts: list[str] = []
            para_fs = 14.0

            for run in para.runs:
                raw = run.text
                if not raw:
                    continue
                text = _esc(raw)

                open_tags: list[str] = []
                close_tags: list[str] = []

                if run.font.bold:
                    open_tags.append("<b>");    close_tags.insert(0, "</b>")
                if run.font.italic:
                    open_tags.append("<i>");    close_tags.insert(0, "</i>")
                if run.font.underline:
                    open_tags.append("<u>");    close_tags.insert(0, "</u>")

                fs = 14.0
                if run.font.size:
                    try:
                        fs = max(6.0, min(96.0, float(run.font.size.pt)))
                    except Exception:
                        pass
                para_fs = max(para_fs, fs)

                font_attrs = f' size="{fs:.1f}"'
                # Use _resolve_color_hex — correctly handles theme/scheme colors
                # that .rgb crashes on (AttributeError silently dropped all colors)
                hex_c = _resolve_color_hex(run.font.color, theme_map)
                if hex_c:
                    font_attrs += f' color="{hex_c}"'

                open_tags.append(f"<font{font_attrs}>")
                close_tags.insert(0, "</font>")

                # ── Hyperlink (XML-based — more reliable than run.hyperlink) ──
                run_html = "".join(open_tags) + text + "".join(close_tags)
                href = _get_run_hyperlink(run, slide_part)
                if href:
                    safe = href.replace("&", "&amp;").replace('"', "&quot;")
                    run_html = f'<link href="{safe}" color="blue">{run_html}</link>'

                html_parts.append(run_html)

            html = "".join(html_parts)
            if not html:
                story.append(Spacer(1, para_fs * 0.4))
                continue

            ps = ParagraphStyle(
                "_sp",
                fontSize=para_fs,
                leading=para_fs * 1.2,
                alignment=_rl_align(para.alignment),
                spaceAfter=2,
            )
            try:
                story.append(RLPara(html, ps))
            except Exception:
                plain = _esc("".join(r.text for r in para.runs))
                story.append(RLPara(plain or " ", ps))

        if story and w > 0 and h > 0:
            pad = 4.0
            frame = Frame(
                x + pad, ybot + pad,
                max(1.0, w - 2 * pad), max(1.0, h - 2 * pad),
                leftPadding=0, rightPadding=0, topPadding=0, bottomPadding=0,
                showBoundary=0,
            )
            frame.addFromList(story, c)

    for slide in prs.slides:
        # Build per-slide theme color map (used by _draw_shape for text colors)
        theme_map = _build_theme_color_map(slide)
        slide_part = slide.part

        # Draw background (white if no theme fill can be resolved)
        drawn_bg = False
        try:
            bg_fill = slide.background.fill
            if bg_fill.type is not None:
                hex_c = _resolve_color_hex(bg_fill.fore_color, theme_map)
                if hex_c:
                    r, g, b = int(hex_c[1:3], 16), int(hex_c[3:5], 16), int(hex_c[5:7], 16)
                    c.setFillColorRGB(r / 255.0, g / 255.0, b / 255.0)
                    c.rect(0, 0, slide_w, slide_h, fill=1, stroke=0)
                    drawn_bg = True
        except Exception:
            pass
        if not drawn_bg:
            c.setFillColorRGB(1, 1, 1)
            c.rect(0, 0, slide_w, slide_h, fill=1, stroke=0)

        for shape in slide.shapes:
            _draw_shape(shape, theme_map, slide_part)

        c.showPage()

    c.save()
    return buffer.getvalue(), len(prs.slides)


# ─────────────────────────────────────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────────────────────────────────────

def convert_docx_to_pdf(file_path: str) -> tuple[bytes, int]:
    """
    Convert a .docx file to PDF bytes.

    Tries converters in order: Word COM → LibreOffice → ReportLab.
    Returns (pdf_bytes, page_count).
    """
    # 1. Microsoft Word (exact "Save As PDF" quality)
    pdf_bytes = _try_word_com(file_path)
    if pdf_bytes:
        logger.info("DOCX converted via Microsoft Word COM")
        return pdf_bytes, _count_pdf_pages(pdf_bytes)

    # 2. LibreOffice (near-identical quality)
    pdf_bytes = _try_libreoffice(file_path)
    if pdf_bytes:
        logger.info("DOCX converted via LibreOffice")
        return pdf_bytes, _count_pdf_pages(pdf_bytes)

    # 3. ReportLab pure-Python fallback
    logger.info("DOCX converted via ReportLab fallback")
    return _docx_reportlab(file_path)


def convert_pptx_to_pdf(file_path: str) -> tuple[bytes, int]:
    """
    Convert a .pptx file to PDF bytes (one page per slide).

    Tries converters in order: PowerPoint COM → LibreOffice → ReportLab.
    Returns (pdf_bytes, page_count).
    """
    # 1. Microsoft PowerPoint (exact "Save As PDF" quality)
    pdf_bytes = _try_powerpoint_com(file_path)
    if pdf_bytes:
        logger.info("PPTX converted via Microsoft PowerPoint COM")
        return pdf_bytes, _count_pdf_pages(pdf_bytes)

    # 2. LibreOffice
    pdf_bytes = _try_libreoffice(file_path)
    if pdf_bytes:
        logger.info("PPTX converted via LibreOffice")
        return pdf_bytes, _count_pdf_pages(pdf_bytes)

    # 3. ReportLab pure-Python fallback
    logger.warning("PPTX COM and LibreOffice both unavailable — using ReportLab fallback (reduced quality)")
    return _pptx_reportlab(file_path)
